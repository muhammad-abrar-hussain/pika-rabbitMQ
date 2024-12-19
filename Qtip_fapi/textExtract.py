"""
Text Extraction Module

This module provides functions to extract text from various document formats, including PDF, PowerPoint, Word, and ODP (OpenDocument Presentation) files.

Features:
- Extract text from PDF files using `PyPDF2`.
- Extract text from PowerPoint files using `python-pptx`.
- Extract text from Word documents (DOCX) using `python-docx`.
- Extract text from ODP files using `odf` libraries.
- Generic function `extract_text` to handle different file types dynamically.

Attributes:
    SUPPORTED_EXTENSIONS (list): List of supported file extensions.
"""

import PyPDF2
from pptx import Presentation
from pathlib import Path
from docx import Document
from odf.opendocument import load
from odf.text import P
import os
# from sklearn.metrics.pairwise import cosine_similarity
# from sklearn.feature_extraction.text import TfidfVectorizer


def extract_text_from_pdf(file_path):
    """
        Extracts text from a PDF file.

        Args:
            file_path (str): The path to the PDF file.

        Returns:
            str: The extracted text from the PDF file.

        Notes:
            - Uses `PyPDF2` to read PDF files.
            - Handles errors gracefully and prints error messages.
        """

    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text


def extract_text_from_pptx(file_path):
    """
        Extracts text from a PowerPoint (PPTX) file.

        Args:
            file_path (str): The path to the PowerPoint file.

        Returns:
            str: The extracted text from the PowerPoint file.

        Notes:
            - Uses `python-pptx` to extract text from slides and their shapes.
            - Handles errors gracefully and prints error messages.
        """

    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
    except Exception as e:
        print(f"Error reading PowerPoint: {e}")
    return text


def extract_text_from_docx(file_path):
    """
        Extracts text from a Word document (DOCX).

        Args:
            file_path (str): The path to the Word document.

        Returns:
            str: The extracted text from the Word document, with each paragraph separated by a newline.

        Notes:
            - Uses the `python-docx` library to process DOCX files.
            - Iterates through all paragraphs in the document and joins their text content with newline characters.
            - Only extracts visible text content; metadata or embedded objects are not included.
        """

    doc = Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])


def extract_text_from_odp(file_path):
    """
        Extracts text from an OpenDocument Presentation (ODP) file.

        Args:
            file_path (str): The path to the ODP file.

        Returns:
            str: The extracted text from the ODP file.

        Notes:
            - Uses `odf.opendocument` to load ODP files.
            - Extracts text content from paragraphs and spans.
            - Handles errors gracefully and returns an error message if extraction fails.
        """

    try:
        doc = load(file_path)
        text = []
        
        # Loop through all paragraphs in the document
        for paragraph in doc.getElementsByType(P):
            paragraph_text = ""
            # Loop through child nodes (e.g., spans, text nodes)
            for node in paragraph.childNodes:
                if node.nodeType == 3:  # Text Node
                    paragraph_text += str(node.data).strip()
                elif node.nodeType == 1 and node.tagName == "text:span":  # Span Node
                    paragraph_text += "".join([str(child.data) for child in node.childNodes if child.nodeType == 3])
            if paragraph_text.strip():  # Avoid adding empty lines
                text.append(paragraph_text)
        
        return "\n".join(text)
    except Exception as e:
        return f"Error reading ODP file: {e}"


def extract_text(file_path):
    """
        Extracts text from a file based on its extension.

        Args:
            file_path (str): The path to the file.

        Returns:
            str: The extracted text from the file.

        Raises:
            ValueError: If the file type is unsupported.

        Notes:
            - Automatically determines the file type based on its extension.
            - Supports DOCX, PDF, PPTX, and ODP formats.
        """

    ext = os.path.splitext(file_path)[1].lower()
    if ext in ['.doc', '.docx']:
        return extract_text_from_docx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.ppt', '.pptx']:
        return extract_text_from_pptx(file_path)
    elif ext == '.odp':
        return extract_text_from_odp(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# file_path = "./assets/test.pdf"
# file_text = extract_text(file_path)
# print(file_text)
#issue-->doc, ppt




# summarized_topics = {
#     "topic_1": "Introduction to Machine Learning",
#     "topic_2": "Applications of Artificial Intelligence",
#     "topic_3": "Basics of Neural Networks",
# }
#
# student_questions = {
#     "q_1": "What are the basics of neural networks?",
#     "q_2": "How is AI applied in healthcare?",
#     "q_3": "what is scope of AI?",
# }
#
# vectorizer = TfidfVectorizer()
#
# topic_vectors = vectorizer.fit_transform(summarized_topics.values())
#
# for question_id, question in student_questions.items():
#     print(f"\nProcessing Question ID: {question_id}")
#
#     question_vector = vectorizer.transform([question])
#
#     similarity_scores = cosine_similarity(question_vector, topic_vectors)
#
#     most_relevant_index = similarity_scores.argmax()
#     relevancy_score = similarity_scores[0, most_relevant_index]
#
#     most_relevant_topic_id = list(summarized_topics.keys())[most_relevant_index]
#
#     label = "Relevant" if relevancy_score > 0.5 else "Irrelevant"
#
#     print(f"Most Relevant Topic ID: {most_relevant_topic_id}")
#     print(f"Relevancy Index: {relevancy_score:.2f}")
#     print(f"Label: {label}")
















































# def clean_and_chunk_text(text, max_token_length=512):
#     cleaned_text = re.sub(r'\s+', ' ', text)  
#     cleaned_text = re.sub(r'[^\w\s]', '', cleaned_text) 

#     # Step 2: Chunk the text
#     words = cleaned_text.split()
#     chunks = [" ".join(words[i:i + max_token_length]) for i in range(0, len(words), max_token_length)]
#     return chunks


# chunks = clean_and_chunk_text(file_text)
# print(chunks)
# for idx, chunk in enumerate(chunks, 1):
#     print(f"Chunk {idx}: {chunk[:512]}")